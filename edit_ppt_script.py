import sys
from pptx import Presentation

def update_slide(presentation, slide_index, old_text, new_text):
    slide = presentation.slides[slide_index]

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)

if __name__ == "__main__":
    old_text_list = sys.argv[1].split('|')
    new_text_list = sys.argv[2].split('|')
    slide_indexes = list(map(int, sys.argv[3].split('|')))
    pptx_path = sys.argv[4]

    # Load the original presentation
    prs = Presentation(pptx_path)

    for old_text, new_text, slide_index in zip(old_text_list, new_text_list, slide_indexes):
        update_slide(prs, slide_index, old_text, new_text)

    # Save the modified presentation
    modified_pptx_path = 'modified_' + pptx_path
    prs.save(modified_pptx_path)

    print("Slides updated.")
