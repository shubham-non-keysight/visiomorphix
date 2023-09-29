import sys
from pptx import Presentation

def update_slide(presentation, slide_index, old_text, new_text):
    slide = presentation.slides[slide_index]

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)

def main():
    input_file = sys.argv[1]
    old_text_list = sys.argv[2].split('|')
    new_text_list = sys.argv[3].split('|')
    slide_indexes = list(map(int, sys.argv[4].split('|')))

    # Load the original presentation
    prs = Presentation(input_file)

    for old_text, new_text, slide_index in zip(old_text_list, new_text_list, slide_indexes):
        update_slide(prs, slide_index, old_text, new_text)

    # Save the modified presentation
    modified_pptx_path = 'modified_presentation.pptx'
    prs.save(modified_pptx_path)

    print("Slides updated.")

if __name__ == "__main__":
    main()
