from pptx import Presentation
import pyautogui
import os
import subprocess

# Function to replace text in a slide while retaining animations
def replace_text_with_animation(slide, old_text, new_text):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if old_text in run.text:
                        # Preserve animations by copying the old text's animations to the new text
                        new_run = run._element
                        run._element = run._r = None
                        new_run.text = new_text

# Load the PowerPoint presentation
presentation = Presentation('C:\\Keysight\\Hackathon\Demo2.ppt')

print('starting')

# Specify the slide index and text to be replaced
slide_index = 0  # Change this to the index of the slide containing the text
old_text = 'Demo_Slide'  # Specify the text you want to replace
new_text = 'Introducing Oscilloscope'  # Specify the new text to replace with

# Replace text on the specified slide
slide = presentation.slides[slide_index]
replace_text_with_animation(slide, old_text, new_text)

# Save the modified presentation
modified_pptx_file = 'modified_presentation.pptx'
presentation.save(modified_pptx_file)

# Export the modified presentation to a video (MP4)
output_video_file = 'output_video.mp4'
pyautogui.hotkey('alt', 'f')
pyautogui.press('v')
pyautogui.press('v')
pyautogui.press('enter')
pyautogui.typewrite(modified_pptx_file, interval=0.1)
pyautogui.press('enter')
pyautogui.press('enter')

# Wait for the video export to complete (you can adjust the sleep duration as needed)
import time
time.sleep(10)

# Rename the exported video to the desired filename
os.rename('presentation.mp4', output_video_file)

# Clean up the temporary modified PowerPoint file
os.remove(modified_pptx_file)

# Optional: You can also delete the temporary folder created by PowerPoint during the export
# Make sure to replace 'C:\\Users\\your_username\\AppData\\Local\\Temp\\ppt_save\\'
# with the actual path to the temporary folder if needed.
temp_folder = 'C:\\Users\\your_username\\AppData\\Local\\Temp\\ppt_save\\'
for filename in os.listdir(temp_folder):
    file_path = os.path.join(temp_folder, filename)
    try:
        if os.path.isfile(file_path):
            os.unlink(file_path)
    except Exception as e:
        print(f"Error deleting {file_path}: {e}")

# Print a message when the process is complete
print(f'Video saved as {output_video_file}')

# Optional: Convert the video to a different format or codec using FFmpeg if needed
# Make sure you have FFmpeg installed and added to your system's PATH.
# For example, you can use the subprocess module to convert to a different format:
# command = f'ffmpeg -i {output_video_file} -c:v libx264 -c:a aac output.mp4'
# subprocess.call(command, shell=True)
