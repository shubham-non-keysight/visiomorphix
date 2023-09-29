import sys
from pptx import Presentation

def update_slide(presentation, slide_index, old_text, new_text):
    slide = presentation.slides[slide_index]

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)

def generate_presentation(input_file):
    prs = Presentation(input_file)
    modified_prs = Presentation()

    for slide in prs.slides:
        new_slide = modified_prs.slides.add_slide(slide.slide_layout)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run_text = run.text
                        new_run = new_slide.shapes[shape.index].text_frame.paragraphs[paragraph.index].runs[0]
                        new_run.text = run_text

    modified_pptx_path = 'temp.pptx'
    modified_prs.save(modified_pptx_path)
    print(f"Temporary modified presentation saved as: {modified_pptx_path}")

if __name__ == "__main__":
    input_file = sys.argv[1]
    generate_presentation(input_file)

    old_text_list = sys.argv[2].split('|')
    new_text_list = sys.argv[3].split('|')
    slide_indexes = list(map(int, sys.argv[4].split('|')))
    pptx_path = 'temp.pptx'

    # Load the original presentation
    prs = Presentation(pptx_path)

    for old_text, new_text, slide_index in zip(old_text_list, new_text_list, slide_indexes):
        update_slide(prs, slide_index, old_text, new_text)

    # Save the modified presentation
    modified_pptx_path = 'modified_presentation.pptx'
    prs.save(modified_pptx_path)

    print("Slides updated.")
